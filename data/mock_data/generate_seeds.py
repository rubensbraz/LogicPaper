import asyncio
import logging
import os
import shutil
import zipfile
from datetime import date, datetime, timedelta
from typing import Any, Dict, List, Optional, Tuple, Union

import xlsxwriter
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt


# --- Configuration ---
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s"
)
logger = logging.getLogger(__name__)

# --- Constants ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_TEMP_DIR = os.path.join(BASE_DIR, ".temp_assets_gen")

# Corporate Colors (RGB)
COLOR_PRIMARY = (31, 78, 120)  # Navy Blue
COLOR_SECONDARY = (231, 76, 60)  # Accent Red
COLOR_TEXT = (51, 51, 51)  # Dark Gray


class SeedGenerator:
    """
    Seed Generator.
    Orchestrates the creation of robust mock data, documentation, assets, and
    high-fidelity templates for system integration testing.
    """

    def __init__(self):
        """Initialize the generator configuration."""
        self.excel_filename = "mock_data.xlsx"
        self.assets_zip_name = "assets.zip"

        # Templates
        self.tmpl_word_brief = "template_brief.docx"
        self.tmpl_word_contract = "template_contract.docx"
        self.tmpl_word_compat = "template_compatibility_word.docx"

        self.tmpl_ppt_pres = "template_presentation.pptx"
        self.tmpl_ppt_compat = "template_compatibility_ppt.pptx"

    # --- Helper Methods ---

    def _ensure_temp_directory(self) -> None:
        """Creates the temporary assets directory safely."""
        os.makedirs(ASSETS_TEMP_DIR, exist_ok=True)

    def _set_no_proof(self, run: Any) -> None:
        """
        Marks a Word run as 'no proofing'.
        Prevents Word spellcheck from breaking Jinja2 tags.

        Args:
            run (Any): The docx Run object.
        """
        try:
            rPr = run._element.get_or_add_rPr()
            noProof = OxmlElement("w:noProof")
            rPr.append(noProof)
        except Exception as e:
            logger.warning(f"Could not set noProof on run: {e}")

    def _style_header_run(
        self, run: Any, size: int = 14, color: Tuple[int, int, int] = COLOR_PRIMARY
    ) -> None:
        """Applies styling to a text run."""
        run.bold = True
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(*color)
        run.font.name = "Calibri"

    # --- Asset Generation ---

    async def generate_assets(self) -> None:
        """
        Generates diverse dummy images and compresses them into a ZIP file.
        """
        logger.info("🎨 Generating High-Res Assets...")
        self._ensure_temp_directory()

        tasks = [
            self._create_image(
                "profile_CEO.jpg", (500, 500), (44, 62, 80), "CEO", (255, 255, 255)
            ),
            self._create_image(
                "profile_CTO.jpg", (500, 500), (39, 174, 96), "CTO", (255, 255, 255)
            ),
            self._create_image(
                "logo_acme.png", (800, 300), (255, 255, 255), "ACME CORP", (31, 78, 120)
            ),
            self._create_image(
                "sig_valid.png",
                (400, 150),
                (255, 255, 255),
                "Signed: J.Doe",
                (0, 0, 128),
            ),
            self._create_image(
                "chart_growth.jpg",
                (800, 400),
                (236, 240, 241),
                "GROWTH CHART",
                (127, 140, 141),
            ),
        ]

        await asyncio.gather(*tasks)
        await asyncio.to_thread(self._zip_assets)

    async def _create_image(
        self,
        filename: str,
        size: Tuple[int, int],
        color: Tuple[int, int, int],
        text: str,
        text_color: Tuple[int, int, int],
    ) -> None:
        """Creates a single dummy image with text."""

        def _blocking_create():
            try:
                img = Image.new("RGB", size, color)
                draw = ImageDraw.Draw(img)
                # Border
                draw.rectangle(
                    [(0, 0), (size[0] - 1, size[1] - 1)], outline=text_color, width=5
                )
                # Text
                try:
                    font = ImageFont.truetype("arial.ttf", 40)
                except IOError:
                    font = ImageFont.load_default()

                # Simple center logic
                bbox = draw.textbbox((0, 0), text, font=font)
                w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
                draw.text(
                    ((size[0] - w) / 2, (size[1] - h) / 2),
                    text,
                    fill=text_color,
                    font=font,
                )

                path = os.path.join(ASSETS_TEMP_DIR, filename)
                img.save(path)
            except Exception as e:
                logger.error(f"Error creating image {filename}: {e}")

        await asyncio.to_thread(_blocking_create)

    def _zip_assets(self) -> None:
        """Compresses temp assets to ZIP and cleans up."""
        zip_path = os.path.join(BASE_DIR, self.assets_zip_name)
        try:
            with zipfile.ZipFile(zip_path, "w") as zipf:
                for root, _, files in os.walk(ASSETS_TEMP_DIR):
                    for file in files:
                        file_path = os.path.join(root, file)
                        zipf.write(
                            file_path, os.path.relpath(file_path, ASSETS_TEMP_DIR)
                        )
            logger.info(f"✅ Assets Library created: {zip_path}")
        except Exception as e:
            logger.error(f"Error zipping assets: {e}")
        finally:
            if os.path.exists(ASSETS_TEMP_DIR):
                shutil.rmtree(ASSETS_TEMP_DIR)

    # --- Excel Generation ---

    async def generate_excel_data(self) -> None:
        """Creates the Excel file with 5 robust rows."""
        logger.info("📊 Generating Excel Data (5 Rows)...")
        filepath = os.path.join(BASE_DIR, self.excel_filename)
        await asyncio.to_thread(self._write_excel_file, filepath)

    def _write_excel_file(self, filepath: str) -> None:
        try:
            wb = xlsxwriter.Workbook(filepath)

            # Formats
            fmt_head = wb.add_format(
                {
                    "bold": True,
                    "font_color": "white",
                    "bg_color": "#2C3E50",
                    "border": 1,
                }
            )
            fmt_date = wb.add_format({"num_format": "yyyy-mm-dd", "border": 1})
            fmt_curr = wb.add_format({"num_format": "$#,##0.00", "border": 1})
            fmt_text = wb.add_format({"border": 1, "text_wrap": True})

            ws = wb.add_worksheet("Data")

            headers = [
                "id",
                "company",
                "ceo_name",
                "founded_date",
                "revenue_q4",
                "growth_pct",
                "is_public",
                "compliance_check",
                "risk_score",
                "email_contact",
                "card_token",
                "logo_img",
                "ceo_img",
                "auth_sig",
            ]
            ws.write_row(0, 0, headers, fmt_head)

            # 5 Robust Rows
            data = [
                # 1. Happy Path (Standard Enterprise)
                [
                    "CORP-001",
                    "Acme Solutions Inc.",
                    "John Sterling",
                    date(1998, 5, 12),
                    1500000.00,
                    0.125,
                    True,
                    False,
                    10,
                    "contact@acme.com",
                    "4532123456789012",
                    "logo_acme.png",
                    "profile_CEO.jpg",
                    "sig_valid.png",
                ],
                # 2. International / Localization (European Context)
                [
                    "CORP-002",
                    "Omega G.m.b.H",
                    "Klaus Müller",
                    date(2005, 11, 30),
                    450500.50,
                    0.052,
                    False,
                    True,
                    20,
                    "info@omega.de",
                    "5500123456789012",
                    "None",
                    "profile_CTO.jpg",
                    "sig_valid.png",
                ],
                # 3. High Risk / Negative Growth
                [
                    "CORP-003",
                    "Cyberdyne Sys",
                    "Sarah Connor",
                    date(2024, 1, 1),
                    -50000.00,
                    -0.25,
                    True,
                    True,
                    30,
                    "alert@skynet.net",
                    "371234567890123",
                    "chart_growth.jpg",
                    "None",
                    "None",
                ],
                # 4. Nulls / Defaults Testing (Crucial for robustness)
                [
                    "CORP-004",
                    "Stealth Startup",
                    None,
                    date(2025, 1, 1),
                    0.00,
                    0.0,
                    False,
                    False,
                    99,
                    None,
                    "1234",
                    "None",
                    "None",
                    "None",
                ],
                # 5. Edge Case (Long Strings & Precision)
                [
                    "CORP-005",
                    "The Very Long Company Name Limited Partnership Global Edition",
                    "Alexander Wolfeschlegelsteinhausenbergerdorff",
                    date(1900, 1, 1),
                    9999999.99,
                    1.0,
                    True,
                    True,
                    10,
                    "ceo@verylongdomainnamecompany.com",
                    "4111111111111111",
                    "logo_acme.png",
                    "profile_CEO.jpg",
                    "sig_valid.png",
                ],
            ]

            for r, row in enumerate(data, 1):
                for c, val in enumerate(row):
                    if isinstance(val, date):
                        ws.write(r, c, val, fmt_date)
                    elif isinstance(val, float):
                        ws.write(r, c, val, fmt_curr)
                    else:
                        ws.write(r, c, val, fmt_text)

            # Autosize approximation
            ws.set_column(0, 13, 20)
            wb.close()
            logger.info(f"✅ Excel Data created: {filepath}")

        except Exception as e:
            logger.error(f"Excel generation failed: {e}")
            raise e

    # --- Word Templates ---

    async def generate_word_templates(self) -> None:
        """Orchestrates Word template creation."""
        logger.info("📄 Generating 3 Word Templates...")

        await asyncio.gather(
            asyncio.to_thread(
                self._create_brief, os.path.join(BASE_DIR, self.tmpl_word_brief)
            ),
            asyncio.to_thread(
                self._create_contract, os.path.join(BASE_DIR, self.tmpl_word_contract)
            ),
            asyncio.to_thread(
                self._create_compatibility,
                os.path.join(BASE_DIR, self.tmpl_word_compat),
            ),
        )

    def _create_brief(self, path: str) -> None:
        """
        Creates 'Brief': A professional summary document with layout tables.
        """
        doc = Document()

        # Header
        h1 = doc.add_heading("EXECUTIVE BRIEFING", 0)
        h1.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Profile Table (Layout)
        table = doc.add_table(rows=1, cols=2)
        table.autofit = False
        table.allow_autofit = False
        table.columns[0].width = Cm(4)
        table.columns[1].width = Cm(12)

        # Left: Image
        cell_img = table.cell(0, 0)
        p = cell_img.paragraphs[0]
        self._set_no_proof(p.add_run("{{ ceo_img | format_image('3.5', '3.5') }}"))

        # Right: Info
        cell_info = table.cell(0, 1)
        p = cell_info.paragraphs[0]
        self._style_header_run(p.add_run("{{ company | format_string('upper') }}\n"))
        p.add_run("CEO: ")
        self._set_no_proof(
            p.add_run("{{ ceo_name | format_logic('default', 'Confidential') }}")
        )
        p.add_run("\nEstablished: ")
        self._set_no_proof(p.add_run("{{ founded_date | format_date('year') }}"))

        doc.add_paragraph().add_run().add_break()

        # Financials
        doc.add_heading("Financial Performance", level=1)
        p = doc.add_paragraph("Fiscal Revenue: ")
        run = p.add_run("{{ revenue_q4 | format_currency('USD') }}")
        run.bold = True
        self._set_no_proof(run)

        p = doc.add_paragraph("Growth Trajectory: ")
        run = p.add_run("{{ growth_pct | format_number('percent', '2') }}")
        self._set_no_proof(run)

        # Status Logic
        doc.add_heading("Risk Assessment", level=1)
        p = doc.add_paragraph("Current Status: ")
        run = p.add_run(
            "{{ risk_score | format_logic('10=Low Risk (Approved)', '20=Medium Risk (Review)', '30=High Risk (Audit)', 'default', 'Unknown') }}"
        )
        run.bold = True
        run.font.color.rgb = RGBColor(200, 0, 0)
        self._set_no_proof(run)

        doc.save(path)

    def _create_contract(self, path: str) -> None:
        """
        Creates 'Contract': Legal document style, heavy on text precision.
        """
        doc = Document()
        doc.add_heading("SERVICE LEVEL AGREEMENT", 0)

        p = doc.add_paragraph("THIS AGREEMENT is dated ")
        self._set_no_proof(p.add_run("{{ founded_date | format_date('iso') }}"))
        p.add_run(" and is made between:")

        # Parties
        doc.add_heading("1. The Parties", level=1)
        p = doc.add_paragraph()
        p.add_run("PROVIDER: LogicPaper Systems Ltd.\n").bold = True
        p.add_run("CLIENT: ")
        run = p.add_run("{{ company | format_string('title') }}")
        run.bold = True
        self._set_no_proof(run)
        p.add_run(" (ID: ")
        self._set_no_proof(p.add_run("{{ id }}"))
        p.add_run(")")

        # Terms
        doc.add_heading("2. Consideration", level=1)
        p = doc.add_paragraph("The Client agrees to pay the sum of ")
        run = p.add_run("{{ revenue_q4 | format_currency('EUR') }}")
        run.bold = True
        self._set_no_proof(run)
        p.add_run(" (")
        run = p.add_run("{{ revenue_q4 | format_number('spell_out', 'en') }}")
        run.italic = True
        self._set_no_proof(run)
        p.add_run(" Euros).")

        # Compliance (Boolean)
        doc.add_heading("3. Compliance & Audit", level=1)
        p = doc.add_paragraph()
        p.add_run("[ ")
        self._set_no_proof(p.add_run("{{ compliance_check | format_bool('check') }}"))
        p.add_run(" ] Compliance Verified")

        p = doc.add_paragraph()
        p.add_run("Publicly Traded: ")
        self._set_no_proof(p.add_run("{{ is_public | format_bool('bool', 'Yes', 'No') }}"))

        # Signatures
        doc.add_paragraph().add_run().add_break()
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Signed for LogicPaper:"
        table.rows[0].cells[1].text = "Signed for Client:"

        # Client Sig Image
        cell = table.rows[1].cells[1]
        p = cell.paragraphs[0]
        self._set_no_proof(p.add_run("{{ auth_sig | format_image('4', '2') }}"))

        doc.save(path)

    def _create_compatibility(self, path: str) -> None:
        """
        Creates 'Compatibility': The Ultimate Unit Test.
        Tests EVERY single strategy method.
        """
        doc = Document()
        doc.add_heading("SYSTEM COMPATIBILITY MATRIX", 0)
        doc.add_paragraph("Automated Verification of all Jinja2 Filters.")

        # Helper to make rows
        def add_test_row(tbl, category, filter_name, syntax, result_tag):
            cells = tbl.add_row().cells
            cells[0].text = category
            cells[1].text = filter_name
            cells[2].text = syntax
            p = cells[3].paragraphs[0]
            r = p.add_run(result_tag)
            r.font.color.rgb = RGBColor(0, 100, 0)  # Green for result
            self._set_no_proof(r)

        table = doc.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        hdrs = table.rows[0].cells
        hdrs[0].text = "Category"
        hdrs[1].text = "Filter"
        hdrs[2].text = "Args"
        hdrs[3].text = "Live Result"

        # 1. Strings
        add_test_row(
            table, "String", "upper", "None", "{{ company | format_string('upper') }}"
        )
        add_test_row(
            table, "String", "lower", "None", "{{ company | format_string('lower') }}"
        )
        add_test_row(
            table, "String", "trim", "None", "{{ company | format_string('trim') }}"
        )
        add_test_row(
            table,
            "String",
            "truncate",
            "10",
            "{{ company | format_string('truncate', '10') }}",
        )
        add_test_row(
            table,
            "String",
            "prefix",
            "ID: ",
            "{{ id | format_string('prefix', 'ID: ') }}",
        )

        # 2. Numbers
        add_test_row(
            table, "Number", "int", "None", "{{ revenue_q4 | format_number('int') }}"
        )
        add_test_row(
            table,
            "Number",
            "float",
            "2",
            "{{ revenue_q4 | format_number('float', '2') }}",
        )
        add_test_row(
            table,
            "Number",
            "percent",
            "1",
            "{{ growth_pct | format_number('percent', '1') }}",
        )
        add_test_row(
            table,
            "Number",
            "spell_out",
            "en",
            "{{ risk_score | format_number('spell_out', 'en') }}",
        )

        # 3. Dates
        add_test_row(
            table, "Date", "iso", "None", "{{ founded_date | format_date('iso') }}"
        )
        add_test_row(
            table, "Date", "long", "None", "{{ founded_date | format_date('long') }}"
        )
        add_test_row(
            table, "Date", "year", "None", "{{ founded_date | format_date('year') }}"
        )
        add_test_row(
            table,
            "Date",
            "add_days",
            "30",
            "{{ founded_date | format_date('add_days', '30') }}",
        )
        add_test_row(
            table,
            "Date",
            "fmt",
            "%d/%m/%Y",
            "{{ founded_date | format_date('fmt', '%d/%m/%Y') }}",
        )

        # 4. Logic & Bool
        add_test_row(
            table,
            "Logic",
            "mapping",
            "10=A, 20=B",
            "{{ risk_score | format_logic('10=Approved', '20=Review', 'default', 'Critical') }}",
        )
        add_test_row(
            table,
            "Logic",
            "default",
            "N/A",
            "{{ ceo_name | format_logic('default', 'VACANT POSITION') }}",
        )
        add_test_row(
            table, "Bool", "Yes / No", "None", "{{ is_public | format_bool('bool', 'Yes', 'No') }}"
        )
        add_test_row(
            table,
            "Bool",
            "check",
            "None",
            "{{ compliance_check | format_bool('check') }}",
        )

        # 5. Mask
        add_test_row(
            table, "Mask", "email", "None", "{{ email_contact | format_mask('email') }}"
        )
        add_test_row(
            table,
            "Mask",
            "credit_card",
            "None",
            "{{ card_token | format_mask('credit_card') }}",
        )

        doc.save(path)

    # --- PowerPoint Templates ---

    async def generate_pptx_templates(self) -> None:
        """Orchestrates PowerPoint template creation."""
        logger.info("🖥️  Generating 2 PPTX Templates...")

        await asyncio.gather(
            asyncio.to_thread(
                self._create_ppt_presentation,
                os.path.join(BASE_DIR, self.tmpl_ppt_pres),
            ),
            asyncio.to_thread(
                self._create_ppt_compatibility,
                os.path.join(BASE_DIR, self.tmpl_ppt_compat),
            ),
        )

    def _create_ppt_presentation(self, path: str) -> None:
        """Creates 'Presentation': Visual dashboard."""
        prs = Presentation()

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "{{ company | format_string('upper') }}"
        slide.placeholders[1].text = (
            "Strategic Review | {{ founded_date | format_date('year') }}"
        )

        # Slide 2: Scorecard
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Financial Scorecard"
        tf = slide.placeholders[1].text_frame

        p = tf.add_paragraph()
        p.text = "Q4 Revenue: {{ revenue_q4 | format_currency('USD') }}"
        p.font.size = PptPt(24)

        p = tf.add_paragraph()
        p.text = "Growth: {{ growth_pct | format_number('percent', '2') }}"
        p.font.size = PptPt(24)

        # Logic Mapping for Color (Text-based simulation)
        p = tf.add_paragraph()
        p.text = "Risk Level: {{ risk_score | format_logic('10=LOW (Green)', '20=MED (Yellow)', '30=HIGH (Red)', 'default', 'UNRATED') }}"
        p.font.bold = True

        prs.save(path)

    def _create_ppt_compatibility(self, path: str) -> None:
        """Creates 'Compatibility': Visual unit test for PPTX."""
        prs = Presentation()

        # Slide 1: Logic & Strings
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Text & Logic Tests"
        tf = slide.placeholders[1].text_frame

        tf.text = "Company: {{ company | format_string('upper') }}"
        p = tf.add_paragraph()
        p.text = "Masked Email: {{ email_contact | format_mask('email') }}"
        p = tf.add_paragraph()
        p.text = "Masked Card: {{ card_token | format_mask('credit_card') }}"

        # Slide 2: Booleans
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Boolean Visuals"
        tf = slide.placeholders[1].text_frame

        p = tf.add_paragraph()
        p.text = "Is Public? {{ is_public | format_bool('check') }}"
        p = tf.add_paragraph()
        p.text = "Compliance? [ {{ compliance_check | format_bool('check') }} ]"

        prs.save(path)

    # --- Main Execution ---

    async def run(self) -> None:
        """Executes the full generation pipeline."""
        logger.info("🚀 Starting Seed Generation...")

        # 1. Assets first (required by Excel)
        await self.generate_assets()

        # 2. Parallel Generation of Excel and Templates
        await asyncio.gather(
            self.generate_excel_data(),
            self.generate_word_templates(),
            self.generate_pptx_templates(),
        )

        logger.info("✨ Generation Complete. Files ready for testing.")


if __name__ == "__main__":
    generator = SeedGenerator()
    asyncio.run(generator.run())
