import io
import logging
import re
from typing import Any, Dict, List, Set

import anyio
from docx import Document
from pptx import Presentation

from app.core.ports import StoragePort

# Configure Logging
logger = logging.getLogger(__name__)


class TemplateValidator:
    """Analyzes template files (DOCX/PPTX) using Regex to extract expected Jinja2 variables.

    Decoupled from filesystem via StoragePort.
    """

    def __init__(self, storage: StoragePort):
        """Initialize the TemplateValidator.

        Args:
            storage (StoragePort): Port for file I/O.
        """
        # Regex Explanation:
        # \{\{\s* -> Match opening braces '{{' and optional whitespace
        # ([a-zA-Z0-9_]+) -> Capture Group 1: The Variable Name (alphanumeric + underscore)
        # .*?           -> Non-greedy match of any character (filters, args, spaces)
        # \}\}          -> Match closing braces '}}'
        self.storage = storage
        self.tag_pattern = re.compile(r"\{\{\s*([a-zA-Z0-9_]+).*?\}\}")

    def _extract_from_text(self, text: str) -> Set[str]:
        """Scans the provided text for Jinja2-style variables using regex.

        Args:
            text (str): The text content to search.

        Returns:
            Set[str]: A set of unique variable names found in the text.
        """
        if not text:
            return set()
        return set(self.tag_pattern.findall(text))

    def extract_tags_from_docx(self, file_path: str) -> Set[str]:
        """Extracts Jinja2 tags from a DOCX file via StoragePort.

        Scans paragraphs, tables, headers, and footers for variables.

        Args:
            file_path (str): The absolute path to the DOCX file.

        Returns:
            Set[str]: A set of unique extracted tags. Returns an empty set on failure.
        """
        tags = set()
        try:
            # Read to memory
            content = self.storage.read_binary(file_path)
            doc = Document(io.BytesIO(content))

            # Body Paragraphs
            for paragraph in doc.paragraphs:
                tags.update(self._extract_from_text(paragraph.text))

            # Tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            tags.update(self._extract_from_text(paragraph.text))

            # Headers and Footers
            for section in doc.sections:
                # Headers
                for header in [
                    section.header,
                    section.first_page_header,
                    section.even_page_header,
                ]:
                    if header:
                        for paragraph in header.paragraphs:
                            tags.update(self._extract_from_text(paragraph.text))
                        for table in header.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    for paragraph in cell.paragraphs:
                                        tags.update(
                                            self._extract_from_text(paragraph.text)
                                        )
                # Footers
                for footer in [
                    section.footer,
                    section.first_page_footer,
                    section.even_page_footer,
                ]:
                    if footer:
                        for paragraph in footer.paragraphs:
                            tags.update(self._extract_from_text(paragraph.text))
                        for table in footer.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    for paragraph in cell.paragraphs:
                                        tags.update(
                                            self._extract_from_text(paragraph.text)
                                        )

            return tags
        except Exception as e:
            logger.error(f"Failed to parse DOCX {file_path}: {e}")
            return set()

    def extract_tags_from_text_file(self, file_path: str) -> Set[str]:
        """Extracts tags from Plain Text or Markdown files via StoragePort.

        Args:
            file_path (str): Path to the text/markdown file.

        Returns:
            Set[str]: A set of extracted tags.
        """
        try:
            content = self.storage.read_text(file_path)
            return self._extract_from_text(content)
        except Exception as e:
            logger.error(f"Failed to parse Text File {file_path}: {e}")
            return set()

    def extract_tags_from_pptx(self, file_path: str) -> Set[str]:
        """Extracts Jinja2 tags from a PPTX file via StoragePort.

        Iterates through slides, shapes, text frames, and tables to find variables.

        Args:
            file_path (str): The absolute path to the PPTX file.

        Returns:
            Set[str]: A set of unique extracted tags. Returns an empty set on failure.
        """
        tags = set()
        try:
            content = self.storage.read_binary(file_path)
            prs = Presentation(io.BytesIO(content))

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Text Frames
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            # We join runs to handle cases where formatting splits the tag
                            full_text = "".join([run.text for run in paragraph.runs])
                            tags.update(self._extract_from_text(full_text))

                            # Fallback: Check raw paragraph text if runs failed to join correctly
                            if not tags:
                                tags.update(self._extract_from_text(paragraph.text))

                    # Tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if hasattr(cell, "text_frame") and cell.text_frame:
                                    tags.update(
                                        self._extract_from_text(cell.text_frame.text)
                                    )

            return tags
        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path}: {e}")
            return set()

    def extract_image_placeholders(self, file_path: str) -> List[Set[str]]:
        """Extracts grouped shape identifiers (Name + Alt Text) from PPTX.

        Args:
            file_path (str): Absolute path to the PPTX file.

        Returns:
            List[Set[str]]: A list where each element is a set of aliases for a single shape.
        """
        shape_groups = []
        try:
            content = self.storage.read_binary(file_path)
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    aliases = set()
                    # 1. Object Name
                    if shape.name:
                        aliases.add(shape.name.strip().lower())

                    # 2. Alt Text (Title/Description)
                    try:
                        if hasattr(shape._element, "nvPicPr"):
                            cNvPr = shape._element.nvPicPr.cNvPr
                        elif hasattr(shape._element, "nvSpPr"):
                            cNvPr = shape._element.nvSpPr.cNvPr
                        else:
                            cNvPr = None

                        if cNvPr is not None:
                            title = cNvPr.get("title")
                            descr = cNvPr.get("descr")
                            if title:
                                aliases.add(title.strip().lower())
                            if descr:
                                aliases.add(descr.strip().lower())
                    except Exception:
                        pass

                    if aliases:
                        shape_groups.append(aliases)
            return shape_groups
        except Exception as e:
            logger.error(f"Failed to extract shapes from PPTX {file_path}: {e}")
            return []

    def extract_image_placeholders_from_docx(self, file_path: str) -> List[Set[str]]:
        """Extracts grouped shape identifiers (Name + Alt Text) from DOCX.

        Args:
            file_path (str): The absolute path to the DOCX file.

        Returns:
            List[Set[str]]: A list where each element is a set of aliases for a single shape.
        """
        shape_groups = []
        try:
            content = self.storage.read_binary(file_path)
            doc = Document(io.BytesIO(content))

            for shape in doc.inline_shapes:
                aliases = set()
                try:
                    doc_props = shape._inline.docPr

                    title = doc_props.get("title")
                    description = doc_props.get("descr")
                    name = doc_props.get("name")

                    if title:
                        aliases.add(title.strip().lower())
                    if description:
                        aliases.add(description.strip().lower())
                    if name:
                        aliases.add(name.strip().lower())

                except AttributeError:
                    continue
                except Exception:
                    continue

                if aliases:
                    shape_groups.append(aliases)

            return shape_groups
        except Exception as e:
            logger.error(f"Failed to extract shapes from DOCX {file_path}: {e}")
            return []

    async def compare(
        self,
        excel_headers: List[str],
        templates_map: Dict[str, str],
        assets_path: str = None,
        assets_error: str = None,
    ) -> Dict[str, Any]:
        """Compares Excel headers and Assets against requirements.

        Args:
            excel_headers (List[str]): List of headers from the Excel file.
            templates_map (Dict[str, str]): Map of filename to file path for templates.
            assets_path (str, optional): Path to the extracted assets directory.
            assets_error (str, optional): Error message if assets extraction failed.

        Returns:
            Dict[str, Any]: Validation report.
        """

        def _blocking_compare():
            # Normalize headers
            available_vars = set(str(h).strip() for h in excel_headers)

            available_assets = set()
            if assets_path and self.storage.exists(assets_path):
                # Search recursively for all files
                files = self.storage.list_files(assets_path, recursive=True)
                for f in files:
                    # We use basename to flatten the structure (e.g. assets/img.png -> img)
                    f_name = self.storage.splitext(self.storage.basename(f))[0].lower()
                    available_assets.add(f_name)

            validation_report = []

            for filename, path in templates_map.items():
                ext = self.storage.splitext(filename)[1].lower()
                required_vars = set()
                required_assets = []

                if ext == ".docx":
                    required_vars = self.extract_tags_from_docx(path)
                    required_assets = self.extract_image_placeholders_from_docx(path)
                elif ext in [".md", ".txt"]:
                    required_vars = self.extract_tags_from_text_file(path)
                elif ext == ".pptx":
                    required_vars = self.extract_tags_from_pptx(path)
                    required_assets = self.extract_image_placeholders(path)

                # Check Variables
                missing_vars = required_vars - available_vars

                # Check Assets
                matched_assets_set = set()
                potential_missing_assets_set = set()

                # Filter noise from unmatched shapes (basic heuristic)
                ignored_patterns = [
                    "title",
                    "subtitle",
                    "footer",
                    "date",
                    "slide number",
                    "placeholder",
                    "rectangle",
                    "textbox",
                    "group",
                    "line",
                    "oval",
                    "logo",
                    "image",
                    "picture",
                ]

                for group in required_assets:
                    # Check if ANY alias in this shape group matches an available asset
                    match_found = False
                    for alias in group:
                        if alias in available_assets:
                            matched_assets_set.add(alias)
                            match_found = True
                            break

                    if not match_found:
                        # If no match, check if this shape is "Interesting" enough to report missing
                        # We pick the "best" alias (not noise) to report
                        best_reporting_name = None
                        for alias in group:
                            is_noise = any(p in alias for p in ignored_patterns)
                            is_digit = alias.isdigit()
                            is_short = len(alias) <= 2

                            if not is_noise and not is_digit and not is_short:
                                best_reporting_name = alias
                                break

                        if best_reporting_name:
                            potential_missing_assets_set.add(best_reporting_name)

                matched_assets = list(matched_assets_set)
                potential_missing_assets = list(potential_missing_assets_set)

                status = "OK"
                if missing_vars:
                    status = "Missing Data"

                if potential_missing_assets:
                    if status == "OK":
                        status = "Warning: Missing Images"

                validation_report.append(
                    {
                        "template": filename,
                        "status": status,
                        "missing_vars": list(missing_vars),
                        "matched_vars": list(
                            required_vars.intersection(available_vars)
                        ),
                        "matched_assets": matched_assets,
                        "potential_missing_assets": potential_missing_assets,
                        "assets_provided": bool(assets_path) and not bool(assets_error),
                        "assets_error": assets_error,
                    }
                )

            # Determine Overall Status
            # Priority: ERROR > WARNING > OK
            final_status = "OK"

            has_error = any(
                item["status"] == "Missing Data" or item.get("assets_error")
                for item in validation_report
            )

            # Check for specific "Warning" string in status, or explicit missing assets
            has_warning = any(
                "Warning" in item["status"]
                or (
                    item.get("potential_missing_assets")
                    and not item.get("assets_error")
                )
                for item in validation_report
            )

            if has_error:
                final_status = "ERROR"
            elif has_warning:
                final_status = "WARNING"

            return {
                "overall_status": final_status,  # Enum-like status: 'OK' | 'WARNING' | 'ERROR'
                "details": validation_report,
            }

        # Offload to thread
        return await anyio.to_thread.run_sync(_blocking_compare)
