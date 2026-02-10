import io
import logging
import re
import zipfile
from typing import Any, Dict, List, Optional

import anyio
from docx.shared import Cm
from docxtpl import DocxTemplate, InlineImage
from jinja2 import BaseLoader, Environment
from PIL import Image
from pptx import Presentation

from app.core.config import settings
from app.core.formatter import DataFormatter
from app.core.ports import ProcessPort, StoragePort

# Configure Logging
logger = logging.getLogger(__name__)


class DocumentEngine:
    """Core engine to manipulate DOCX/PPTX and convert to PDF.

    Template-based formatting via Jinja2 filters and custom PPTX regex parsing.
    Decoupled from filesystem/subprocess via Ports.
    """

    def __init__(self, temp_dir: str, storage: StoragePort, process: ProcessPort):
        """Initialize the Engine.

        Args:
            temp_dir (str): Base directory for temporary files.
            storage (StoragePort): Port for file I/O.
            process (ProcessPort): Port for external processes.
        """
        self.temp_dir = temp_dir
        self.storage = storage
        self.process = process
        self.formatter = DataFormatter()

    def _build_assets_map(self, assets_path: str) -> Dict[str, str]:
        """Builds a map of asset names to absolute paths.

        Maps both:
        1. Filename with extension (lowercase) -> Full Path
        2. Filename without extension (lowercase) -> Full Path

        Args:
            assets_path (str): Path to the assets directory.

        Returns:
            Dict[str, str]: Map of asset identifiers to file paths.
        """
        assets_map = {}
        if not assets_path or not self.storage.exists(assets_path):
            return assets_map

        files = self.storage.list_dir(assets_path)
        for f in files:
            full_path = self.storage.join_path(assets_path, f)

            # Map full filename (lowercase)
            assets_map[f.lower()] = full_path

            # Map stem (lowercase)
            f_name, _ = self.storage.splitext(f)
            # If multiple files have same stem, last one wins (ambiguity accepted)
            assets_map[f_name.lower()] = full_path

        return assets_map

    def _get_image_object(
        self,
        tpl: DocxTemplate,
        value: Any,
        args: List[str],
        assets_map: Dict[str, str],
        assets_path: Optional[str] = None,
    ) -> Any:
        """Generates an InlineImage object for insertion into a DOCX template.

        Parses arguments to determine dimensions and locates the image file within
        the assets directory using the provided map.

        Args:
            tpl (DocxTemplate): The active DocxTemplate instance.
            value (Any): The filename or value indicating the image.
            args (List[str]): Additional arguments (e.g., width_cm, height_cm).
            assets_map (Dict[str, str]): Map of available assets.
            assets_path (Optional[str]): Optional path for logging/fallback. Defaults to None.

        Returns:
            Any: An InlineImage object if successful, or an error string if failed.
        """
        # Resolve dimensions and filename from strategy
        img_data = self.formatter._apply_strategy("image", value, *args)
        filename = img_data.get("filename")

        if not filename or filename == "None":
            return ""

        # Resolve File Path using Map
        filename_lower = filename.lower()
        img_path = assets_map.get(filename_lower)

        if not img_path:
            if assets_path:
                logger.warning(f"Image not found in assets: {filename}")
            return "[IMAGE NOT FOUND]"

        try:
            # Safety check: Verification is only possible if assets_path is provided.
            # Note: assets_map is built from list_dir(assets_path), so entries strictly
            # from that directory are generally safe
            if assets_path and not self.storage.is_safe_path(assets_path, img_path):
                logger.warning(f"Security: Path traversal detected for {filename}")
                return "[Invalid Image Path]"

            # Read image to memory to fully decouple from file path requirements of some libraries
            # InlineImage supports file path or file-like object
            img_bytes = self.storage.read_binary(img_path)
            img_stream = io.BytesIO(img_bytes)

            width = Cm(float(img_data["width"])) if img_data.get("width") else None
            height = Cm(float(img_data["height"])) if img_data.get("height") else None

            return InlineImage(tpl, img_stream, width=width, height=height)
        except Exception as e:
            logger.error(f"Error loading image '{filename}': {e}")
            return f"[IMAGE ERROR: {str(e)}]"

    def _parse_and_replace_pptx_text(self, text: str, context: Dict[str, Any]) -> str:
        """Parses PPTX text for pseudo-Jinja tags and applies formatting manually.

        Args:
            text (str): The raw text containing Jinja2 variables.
            context (Dict[str, Any]): The data context.

        Returns:
            str: The processed text with variables replaced and filters applied.
        """
        # Regex to capture: {{ variable | filter('arg1', 'arg2') }}
        # Group 1: Variable Name
        # Group 2: Full Filter String (optional)
        pattern = re.compile(r"\{\{\s*([a-zA-Z0-9_]+)(\s*\|.*?)?\s*\}\}")

        def replace_match(match):
            """Regex callback to replace Jinja2 tags with formatted values."""
            var_name = match.group(1)
            filter_part = match.group(2)  # e.g., " | format_string('upper')"

            # Get Raw Value
            value = context.get(var_name, "")

            if not filter_part:
                return str(value)

            # Parse Filter Logic
            # Expected format: | filter_name('arg1', 'arg2')
            try:
                # Remove pipe and whitespace
                content = filter_part.strip().lstrip("|").strip()
                # Split filter name from args: "format_string('upper')" -> "format_string", "'upper'"
                if "(" in content and content.endswith(")"):
                    f_name, args_raw = content.split("(", 1)
                    args_raw = args_raw[:-1]  # Remove trailing )
                else:
                    f_name = content
                    args_raw = ""

                args = []
                if args_raw:
                    # Remove quotes and split
                    parts = args_raw.split(",")
                    args = [p.strip().strip("'").strip('"') for p in parts]

                # Map filter name to Strategy Name
                # "format_date" -> "date"
                strategy_map = {
                    "format_string": "string",
                    "format_number": "number",
                    "format_currency": "number",
                    "format_date": "date",
                    "format_bool": "bool",
                    "format_mask": "mask",
                    "format_logic": "logic",
                }

                # Special handling for aliases
                strat_key = strategy_map.get(f_name)
                final_args = args

                if f_name == "format_currency":
                    strat_key = "number"
                    final_args = ["currency"] + args

                if strat_key:
                    return str(
                        self.formatter._apply_strategy(strat_key, value, *final_args)
                    )
                else:
                    # Fallback if unknown filter
                    return str(value)

            except Exception as e:
                logger.error(f"PPTX Filter Error parsing '{filter_part}': {e}")
                return str(value)

        return pattern.sub(replace_match, text)

    def _remove_office_thumbnail_stream(self, content: bytes) -> bytes:
        """Removes the thumbnail from the Office file stream.

        Args:
            content (bytes): The binary content of the Office file.

        Returns:
            bytes: The binary content with the thumbnail removed.
        """
        try:
            in_buffer = io.BytesIO(content)
            out_buffer = io.BytesIO()
            with zipfile.ZipFile(in_buffer, "r") as zin:
                with zipfile.ZipFile(out_buffer, "w") as zout:
                    for item in zin.infolist():
                        if settings.OFFICE_THUMBNAIL_PATH not in item.filename:
                            zout.writestr(item, zin.read(item.filename))
            return out_buffer.getvalue()
        except Exception as e:
            logger.warning(f"Could not strip thumbnail in stream: {e}")
            return content

    async def process_docx(
        self,
        template_path: str,
        output_path: str,
        context: Dict[str, Any],
        assets_path: Optional[str] = None,
    ) -> bool:
        """Renders a DOCX template using Jinja2 context and Custom Filters.

        This method runs in a separate thread to prevent blocking the asyncio event loop.

        Args:
            template_path (str): The absolute path to the template file.
            output_path (str): The absolute path where the rendered file will be saved.
            context (Dict[str, Any]): The data context dictionary for rendering.
            assets_path (Optional[str]): The path to the assets directory. Defaults to None.

        Returns:
            bool: True if rendering was successful, False otherwise.

        Raises:
            Exception: Propagates any exception that occurs during rendering.
        """

        def _blocking_docx_render():
            """Executes blocking DOCX rendering logic."""
            try:
                # Read template to memory
                tpl_bytes = self.storage.read_binary(template_path)
                tpl = DocxTemplate(io.BytesIO(tpl_bytes))

                jinja_env = Environment(autoescape=True)
                filters = self.formatter.get_jinja_filters()
                jinja_env.filters.update(filters)

                # Build Assets Map
                assets_map = {}
                if assets_path:
                    assets_map = self._build_assets_map(assets_path)

                def format_image_wrapper(val, *args):
                    if not assets_map:
                        return "[NO ASSETS FOUND]"
                    return self._get_image_object(
                        tpl, val, list(args), assets_map, assets_path
                    )

                jinja_env.filters["format_image"] = format_image_wrapper

                tpl.render(context, jinja_env=jinja_env)

                if assets_map:
                    self._replace_images_in_docx_shapes(tpl, assets_map)

                # Save to memory
                out_stream = io.BytesIO()
                tpl.save(out_stream)
                out_bytes = out_stream.getvalue()

                # Post-process (Thumbnail removal)
                final_bytes = self._remove_office_thumbnail_stream(out_bytes)

                # Write to storage
                self.storage.write_binary(output_path, final_bytes)
                return True

            except Exception as e:
                logger.error(f"DOCX Render Error: {e}")
                raise e

        return await anyio.to_thread.run_sync(_blocking_docx_render)

    def _replace_images_in_docx_shapes(
        self, tpl: DocxTemplate, assets_map: Dict[str, str]
    ) -> None:
        """Replaces images in DOCX inline shapes based on Alt Text or Object Name.

        Iterates through all inline shapes in the document. Checks for matches in:
        1. Alt Text Title
        2. Alt Text Description
        3. Object Name (Selection Pane name)

        Args:
            tpl (DocxTemplate): The active template object.
            assets_map (Dict[str, str]): A dictionary mapping asset names (lowercase,
                                         no extension) to their absolute file paths.
        """
        try:
            # tpl is a DocxTemplate, but we treat it as a Document for shape iteration
            for shape in tpl.inline_shapes:
                try:
                    # Access Alt Text (Title/Description) via XML attributes
                    # 'docPr' is the CT_NonVisualDrawingProps element
                    doc_props = shape._inline.docPr
                    title = doc_props.get("title", "").strip().lower()
                    description = doc_props.get("descr", "").strip().lower()
                    obj_name = doc_props.get("name", "").strip().lower()

                    match_name = None
                    if title in assets_map:
                        match_name = title
                    elif description in assets_map:
                        match_name = description
                    elif obj_name in assets_map:
                        match_name = obj_name

                    if match_name:
                        img_path = assets_map[match_name]
                        if not self.storage.exists(img_path):
                            logger.warning(
                                f"Asset defined in DOCX shape '{match_name}' not found at: {img_path}"
                            )
                            continue

                        # Read the New Image Data
                        new_img_bytes = self.storage.read_binary(img_path)

                        # Replace the Image Blob
                        # 1. Get Relationship ID from the blip element
                        # path: shape -> inline -> graphic -> graphicData -> pic -> blipFill -> blip -> embed
                        graphic = shape._inline.graphic
                        pic = graphic.graphicData.pic
                        blip = pic.blipFill.blip
                        r_id = blip.embed

                        # 2. Get the related Image Part using the Relationship ID
                        # tpl.part is the Main Document Part
                        if r_id in tpl.part.related_parts:
                            related_part = tpl.part.related_parts[r_id]

                            # 3. Calculate Aspect Ratio Preservation
                            # Get original dimensions (EMU)
                            # shape.width and shape.height are in EMUs (English Metric Units)
                            # 1 inch = 914400 EMUs
                            # 1 cm = 360000 EMUs
                            ph_width = shape.width
                            ph_height = shape.height

                            with Image.open(io.BytesIO(new_img_bytes)) as img:
                                img_w_px, img_h_px = img.size
                                img_ratio = img_w_px / img_h_px

                            ph_ratio = ph_width / ph_height

                            if img_ratio > ph_ratio:
                                # Image is wider relative to placeholder -> Fit to Width
                                new_width = ph_width
                                new_height = int(ph_width / img_ratio)
                            else:
                                # Image is taller relative to placeholder -> Fit to Height
                                new_height = ph_height
                                new_width = int(ph_height * img_ratio)

                            # 4. Update Shape Dimensions
                            shape.width = new_width
                            shape.height = new_height

                            # 5. Overwrite the Blob directly
                            # This updates the underlying image data. Note that if multiple shapes
                            # share the same rId, ALL of them will be updated
                            related_part._blob = new_img_bytes

                            logger.info(
                                f"Replaced DOCX image shape '{match_name}' with asset '{self.storage.basename(img_path)}' (Resized to {new_width}x{new_height} EMU)."
                            )

                except AttributeError:
                    # Shape does not have the expected XML structure (e.g. no docPr)
                    continue
                except Exception:
                    # General safety catch for individual shapes
                    continue

        except Exception as e:
            logger.error(f"DOCX Image Replacement Error: {e}")

    async def process_text(
        self,
        template_path: str,
        output_path: str,
        context: Dict[str, Any],
    ) -> bool:
        """Renders Text-based templates (MD, TXT).

        Args:
            template_path (str): Path to the template file.
            output_path (str): Path to save the rendered file.
            context (Dict[str, Any]): Data context for rendering.

        Returns:
            bool: True if successful, False otherwise.

        Raises:
            Exception: If rendering fails.
        """
        try:
            # Read content
            content = self.storage.read_text(template_path)

            # Render
            jinja_env = Environment(loader=BaseLoader(), autoescape=False)

            # Register Standard Filters (String, Date, Number, etc.)
            filters = self.formatter.get_jinja_filters()
            jinja_env.filters.update(filters)

            # Handle Image Filter for Text
            # In text files, we return the filename string so the user can use it in Markdown tags:
            # Example: ![Alt]({{ photo | format_image }}) -> ![Alt](photo.jpg)
            jinja_env.filters["format_image"] = lambda val, *args: (
                str(val) if val else ""
            )

            # Render
            template = jinja_env.from_string(content)
            rendered_content = template.render(context)

            # Write Output
            self.storage.write_text(output_path, rendered_content)

            return True

        except Exception as e:
            logger.error(f"Text/MD Render Error: {e}")
            raise e

    def _replace_images_in_slide(self, slide: Any, assets_map: Dict[str, str]) -> None:
        """Replaces images in PPTX slides based on Alt Text or Shape Name.

        Iterates through all shapes in the slide. Checks for matches in:
        1. Alt Text Title
        2. Alt Text Description
        3. Shape Name (Selection Pane name)

        Args:
            slide (Any): The PPTX slide object to process.
            assets_map (Dict[str, str]): A dictionary mapping asset names (lowercase,
                                         no extension) to their absolute file paths.
        """
        # Shapes to remove after replacement to avoid modifying list while iterating
        shapes_to_remove = []

        for shape in slide.shapes:
            # Check for matches in multiple properties
            # 1. Object Name
            shape_name_clean = shape.name.strip().lower()

            # 2. Alt Text (Title/Description)
            alt_title = ""
            alt_descr = ""
            try:
                # Accessing XML properties for Alt Text
                # Pictures use nvPicPr, Shapes use nvSpPr
                if hasattr(shape._element, "nvPicPr"):
                    cNvPr = shape._element.nvPicPr.cNvPr
                elif hasattr(shape._element, "nvSpPr"):
                    cNvPr = shape._element.nvSpPr.cNvPr
                else:
                    cNvPr = None

                if cNvPr is not None:
                    alt_title = cNvPr.get("title", "").strip().lower()
                    alt_descr = cNvPr.get("descr", "").strip().lower()
            except Exception:
                pass

            # Prioritized matching
            match_name = None
            if alt_title in assets_map:
                match_name = alt_title
            elif alt_descr in assets_map:
                match_name = alt_descr
            elif shape_name_clean in assets_map:
                match_name = shape_name_clean

            if match_name:
                try:
                    img_path = assets_map[match_name]

                    # Get placeholder dimensions & position
                    ph_left = shape.left
                    ph_top = shape.top
                    ph_width = shape.width
                    ph_height = shape.height

                    # Load Image Bytes via Storage
                    if not self.storage.exists(img_path):
                        logger.warning(f"Image asset not found: {img_path}")
                        continue

                    img_bytes = self.storage.read_binary(img_path)

                    # Calculate Aspect Ratio Preservation
                    with Image.open(io.BytesIO(img_bytes)) as img:
                        img_w_px, img_h_px = img.size
                        img_ratio = img_w_px / img_h_px

                    ph_ratio = ph_width / ph_height

                    if img_ratio > ph_ratio:
                        # Image is wider relative to placeholder -> Fit to Width
                        new_width = ph_width
                        new_height = int(ph_width / img_ratio)
                    else:
                        # Image is taller relative to placeholder -> Fit to Height
                        new_height = ph_height
                        new_width = int(ph_height * img_ratio)

                    # Center the image
                    left_offset = (ph_width - new_width) // 2
                    top_offset = (ph_height - new_height) // 2

                    new_left = ph_left + left_offset
                    new_top = ph_top + top_offset

                    # Insert Picture
                    slide.shapes.add_picture(
                        io.BytesIO(img_bytes), new_left, new_top, new_width, new_height
                    )

                    # Mark for removal
                    shapes_to_remove.append(shape)

                except Exception as e:
                    logger.error(
                        f"Failed to replace image for shape '{shape.name}': {e}"
                    )

        # Remove replaced shapes
        # Fallback safe removal for python-pptx
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    async def process_pptx(
        self,
        template_path: str,
        output_path: str,
        context: Dict[str, Any],
        assets_path: Optional[str] = None,
    ) -> bool:
        """Renders a PPTX template using python-pptx.

        This method runs in a separate thread to prevent blocking the asyncio event loop.
        It iterates through slides and replaces placeholders with values from the context.

        Args:
            template_path (str): The absolute path to the template file.
            output_path (str): The absolute path where the rendered file will be saved.
            context (Dict[str, Any]): The data context dictionary for rendering.
            assets_path (Optional[str]): The path to the assets directory. Defaults to None.

        Returns:
            bool: True if rendering was successful, False otherwise.

        Raises:
            Exception: Propagates any exception that occurs during rendering.
        """

        def _blocking_pptx_render():
            try:
                # Read template to memory
                tpl_bytes = self.storage.read_binary(template_path)
                prs = Presentation(io.BytesIO(tpl_bytes))

                # Build Assets Map (Name -> Path)
                assets_map = {}
                if assets_path:
                    assets_map = self._build_assets_map(assets_path)

                for slide in prs.slides:
                    for shape in slide.shapes:
                        # Text Frames
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                full_text = "".join(run.text for run in paragraph.runs)
                                if "{{" in full_text:
                                    new_text = self._parse_and_replace_pptx_text(
                                        full_text, context
                                    )
                                    if paragraph.runs:
                                        paragraph.runs[0].text = new_text
                                        for i in range(1, len(paragraph.runs)):
                                            paragraph.runs[i].text = ""

                        # Tables
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text_frame:
                                        for paragraph in cell.text_frame.paragraphs:
                                            full_cell_text = "".join(
                                                run.text for run in paragraph.runs
                                            )
                                            if "{{" in full_cell_text:
                                                new_text = (
                                                    self._parse_and_replace_pptx_text(
                                                        full_cell_text, context
                                                    )
                                                )
                                                if paragraph.runs:
                                                    paragraph.runs[0].text = new_text
                                                    for i in range(
                                                        1, len(paragraph.runs)
                                                    ):
                                                        paragraph.runs[i].text = ""

                    # Image Replacement (if available)
                    if assets_map:
                        self._replace_images_in_slide(slide, assets_map)

                # Save to memory
                out_stream = io.BytesIO()
                prs.save(out_stream)
                out_bytes = out_stream.getvalue()

                # Post-process
                final_bytes = self._remove_office_thumbnail_stream(out_bytes)

                # Write to storage
                self.storage.write_binary(output_path, final_bytes)
                return True
            except Exception as e:
                logger.error(f"PPTX Render Error: {e}")
                raise e

        # Offload to thread
        return await anyio.to_thread.run_sync(_blocking_pptx_render)

    async def convert_to_pdf(self, input_path: str, output_dir: str) -> bool:
        """Converts an Office document to PDF via ProcessPort.

        Args:
            input_path (str): The absolute path to the input document.
            output_dir (str): The directory where the output PDF should be saved.

        Returns:
            bool: True if the conversion process completes successfully.

        Raises:
            Exception: If the conversion fails (non-zero exit code) or times out.
        """
        try:
            cmd = [
                settings.LIBREOFFICE_BINARY,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                output_dir,
                input_path,
            ]

            result = await self.process.run_command(
                cmd, timeout=settings.LIBREOFFICE_TIMEOUT
            )

            if result["returncode"] != 0:
                logger.error(f"LibreOffice failed: {result['stderr'].decode()}")
                raise Exception("PDF Conversion Failed: LibreOffice Error")

            return True

        except Exception as e:
            logger.error(f"PDF Convert Error: {e}")
            raise e
